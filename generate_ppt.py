from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # --- Slide 1: Title ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "CAPSTONE PROJECT\nAGRICARE AI: INTELLIGENT CROP DISEASE DETECTION SYSTEM"
    subtitle.text = "Presented By:\nRohit Prajapati\nJaypee University of Engineering And Technology\nB.Tech-CSE\nSTU67f7493e603021744259390"

    # --- Slide 2: Outline ---
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "OUTLINE"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Problem Statement"
    p = tf.add_paragraph()
    p.text = "Proposed System/Solution"
    p = tf.add_paragraph()
    p.text = "System Development Approach"
    p = tf.add_paragraph()
    p.text = "Algorithm & Deployment"
    p = tf.add_paragraph()
    p.text = "Result"
    p = tf.add_paragraph()
    p.text = "Conclusion"
    p = tf.add_paragraph()
    p.text = "Future Scope"
    p = tf.add_paragraph()
    p.text = "References"

    # --- Slide 3: Problem Statement ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "PROBLEM STATEMENT"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Farmers struggle to quickly and accurately identify crop diseases."
    p = tf.add_paragraph()
    p.text = "Lack of access to agricultural experts leads to delayed treatment."
    p = tf.add_paragraph()
    p.text = "Manual disease identification is time-consuming, expensive, and error-prone."
    p = tf.add_paragraph()
    p.text = "This leads to significant crop yield and quality losses."

    # --- Slide 4: Proposed Solution ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "PROPOSED SOLUTION"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "AgriCare AI: An intelligent web application for instant disease diagnosis."
    p = tf.add_paragraph()
    p.text = "Farmers upload an image of a leaf."
    p = tf.add_paragraph()
    p.text = "An AI model (CNN) analyzes the image and predicts the disease with a confidence score."
    p = tf.add_paragraph()
    p.text = "The system instantly provides actionable treatment and prevention recommendations."

    # --- Slide 5: System Approach ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "SYSTEM APPROACH"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Data Collection: PlantVillage dataset (>50,000 images, 8 classes like Tomato Early Blight, Corn Rust, etc.)"
    p = tf.add_paragraph()
    p.text = "Data Preprocessing: Image resizing to 224x224, pixel normalization (0-1), and data augmentation (rotation, zoom, flipping)."
    p = tf.add_paragraph()
    p.text = "System Requirements: Python, HTML/CSS/JS frontend."
    p = tf.add_paragraph()
    p.text = "Libraries Used: TensorFlow, Keras, OpenCV, Flask, NumPy."

    # --- Slide 6: Algorithm & Deployment ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "ALGORITHM & DEPLOYMENT"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Algorithm Selection: Convolutional Neural Network (CNN) for effective pattern recognition and image classification."
    p = tf.add_paragraph()
    p.text = "Architecture: Input Layer -> Conv2D -> MaxPooling -> Flatten -> Dense -> Dropout -> Output (Softmax)."
    p = tf.add_paragraph()
    p.text = "Training Process: 70% Train, 15% Val, 15% Test. Adam optimizer, Categorical Crossentropy loss."
    p = tf.add_paragraph()
    p.text = "Deployment: Model exported as .h5 and served via a Flask web server, enabling accessible browser-based usage."

    # --- Slide 7: Result ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "RESULT"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Achieved accurate prediction of multiple crop diseases."
    p = tf.add_paragraph()
    p.text = "Prediction Time: < 3 seconds per image."
    p = tf.add_paragraph()
    p.text = "Target Metrics: >90% Accuracy, >88% Precision and Recall."
    p = tf.add_paragraph()
    p.text = "Provides an easy-to-use interface displaying confidence score and treatment plans."

    # --- Slide 8: Conclusion ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "CONCLUSION"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "AgriCare AI successfully provides an accessible tool for early crop disease detection."
    p = tf.add_paragraph()
    p.text = "Directly aligns with UN SDG 2 (Zero Hunger) by preventing crop losses and improving productivity."
    p = tf.add_paragraph()
    p.text = "Demonstrates the practical application of Computer Vision in agriculture."

    # --- Slide 9: Future Scope ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "FUTURE SCOPE"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Development of a mobile application (Android/iOS) for offline field use."
    p = tf.add_paragraph()
    p.text = "Multi-language support and voice assistance for better farmer accessibility."
    p = tf.add_paragraph()
    p.text = "Real-time camera detection feature."
    p = tf.add_paragraph()
    p.text = "Integration with weather APIs for proactive disease prediction."

    # --- Slide 10: References ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "REFERENCES"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "PlantVillage Dataset (Hughes, D., & Salathe, M.)"
    p = tf.add_paragraph()
    p.text = "TensorFlow/Keras Documentation (https://www.tensorflow.org/)"
    p = tf.add_paragraph()
    p.text = "Flask Documentation (https://flask.palletsprojects.com/)"

    # --- Slide 11: Deployment Link / Github Link ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Deployment / GitHub Link"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "GitHub Link: https://github.com/ROHIT-PR22/agricare-ai"
    
    p = tf.add_paragraph()
    p.text = "Deployment: https://agricare-ai-blush.vercel.app/"

    # --- Slide 12: Thank You ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "THANK YOU"
    
    prs.save('AgriCareAI_Capstone_Presentation_Final.pptx')
    print("Presentation generated successfully: AgriCareAI_Capstone_Presentation_Final.pptx")

if __name__ == '__main__':
    create_presentation()
